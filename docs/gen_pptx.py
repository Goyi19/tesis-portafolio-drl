"""
Genera la presentación detallada para el profesor guía.
Resume visualmente el Preprocesamiento, el MDP, y la mecánica interna de PPO y SAC.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paleta de Colores ───────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)  # Slate 900
BG_CARD   = RGBColor(0x1E, 0x29, 0x3B)  # Slate 800
ACCENT    = RGBColor(0x38, 0xBD, 0xF8)  # Sky 400
ACCENT2   = RGBColor(0x10, 0xB9, 0x81)  # Emerald 500
WHITE     = RGBColor(0xF8, 0xFA, 0xFC)
GRAY      = RGBColor(0x94, 0xA3, 0xB8)
LIGHT     = RGBColor(0xE2, 0xE8, 0xF0)
ORANGE    = RGBColor(0xFB, 0x92, 0x3C)
RED       = RGBColor(0xF8, 0x71, 0x71)
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = align
    return tf

def add_para(tf, text, font_size=16, color=WHITE, indent=0, bold=False, font_name="Segoe UI"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.bold = bold
    p.level = indent

def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.width = Pt(1)
    shape.line.color.rgb = GRAY
    return shape

# ════════════════════════════════════════════════════════════════
# SLIDE 1 – PORTADA
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_text(s, 1.5, 2.5, 10.3, 1.2, "Optimización de Portafolios con Deep Reinforcement Learning", font_size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 3.8, 10.3, 0.6, "Metodología, Formulación MDP y Mecánica de PPO / SAC", font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 6.0, 10.3, 0.5, "Presentación de Avance Metodológico al Profesor Guía", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 – PREPROCESAMIENTO (DATA PIPELINE)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.5, 11, 0.6, "1. Preprocesamiento: Transformación del Dataset", font_size=32, bold=True, color=ACCENT)

# Card T=163
add_card(s, 0.8, 1.5, 3.5, 5.2)
tf = add_text(s, 1.1, 1.7, 3, 0.5, "Data Split (T=163)", font_size=20, bold=True, color=ORANGE)
add_para(tf, "Train: 114 Semanas (70%)", font_size=15, color=WHITE, bold=True)
add_para(tf, "Aprendizaje de la política.", font_size=14, color=GRAY)
add_para(tf, "", font_size=6)
add_para(tf, "Val: 24 Semanas (15%)", font_size=15, color=WHITE, bold=True)
add_para(tf, "Ajuste de hiperparámetros.", font_size=14, color=GRAY)
add_para(tf, "", font_size=6)
add_para(tf, "Test: 25 Semanas (15%)", font_size=15, color=WHITE, bold=True)
add_para(tf, "Evaluación Out-of-Sample pura versus GAMS.", font_size=14, color=GRAY)

# Card Momentos
add_card(s, 4.6, 1.5, 7.9, 2.5)
tf = add_text(s, 4.9, 1.7, 7.3, 0.5, "Estimación Rolling Ponderada (Evitar Data Leakage)", font_size=20, bold=True, color=ACCENT2)
add_para(tf, "Se calculan mu (μ) y sigma (σ) por régimen (bear/bull) usando una ventana expansiva, ponderando por la probabilidad HMM observada hasta el periodo t.", font_size=14, color=WHITE)
add_para(tf, "", font_size=6)
add_para(tf, "Las variables Mix:", font_size=16, color=ORANGE, bold=True)
add_para(tf, "Conformamos una única expectativa uniendo ambos regímenes:", font_size=14, color=GRAY)
add_para(tf, "μ_mix = P(bear)*μ_bear + P(bull)*μ_bull", font_size=14, color=ACCENT, font_name="Courier New")
add_para(tf, "σ_mix = P(bear)*σ_bear + P(bull)*σ_bull", font_size=14, color=ACCENT, font_name="Courier New")

# Card Redundancia & Norm
add_card(s, 4.6, 4.2, 7.9, 2.5)
tf = add_text(s, 4.9, 4.4, 7.3, 0.5, "Normalización y Aplanamiento", font_size=20, bold=True, color=PURPLE)
add_para(tf, "Manejo de Probabilidades (Decisión Clave):", font_size=15, color=WHITE, bold=True)
add_para(tf, "Dado que P(bear) y P(bull) ya se inyectan en el cálculo de las variables Mix, NO se pasan crudas al agente para evitar colinealidad y redundancia en la red neuronal.", font_size=13, color=GRAY)
add_para(tf, "Normalización Z-Score Expansiva:", font_size=15, color=WHITE, bold=True)
add_para(tf, "Aplicada a variables mix para evitar gradientes explosivos (Vanishing/Exploding Gradients).", font_size=13, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 – MDP: ESTADO Y OBSERVACIÓN
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.5, 11, 0.6, "2. Proceso de Decisión de Markov (MDP): El Entorno", font_size=32, bold=True, color=ACCENT)

# Estado Interno vs Observacion
add_card(s, 0.8, 1.5, 5.6, 5.2)
tf = add_text(s, 1.1, 1.7, 5, 0.5, "Estado Interno y Observación (o_t)", font_size=20, bold=True, color=ORANGE)
add_para(tf, "Estado Interno s_t = (x_t, w_t-1, t)", font_size=16, color=WHITE, font_name="Courier New", bold=True)
add_para(tf, "Capital actual, pesos anteriores y tiempo.", font_size=14, color=GRAY)
add_para(tf, "", font_size=6)
add_para(tf, "Vector de Observación o_t:", font_size=16, color=WHITE, bold=True)
add_para(tf, "Lo que ve la red neuronal:", font_size=14, color=GRAY)
add_para(tf, "• Capital normalizado (x_t / x_0)", font_size=14, color=ACCENT)
add_para(tf, "• Pesos de periodo anterior (w_SPX, w_CMC)", font_size=14, color=ACCENT)
add_para(tf, "• Tiempo normalizado (t / T)", font_size=14, color=ACCENT)
add_para(tf, "• μ_mix, σ_mix estandarizados (Features)", font_size=14, color=ACCENT)

# Accion y Transicion
add_card(s, 6.7, 1.5, 5.8, 5.2)
tf = add_text(s, 7.0, 1.7, 5.2, 0.5, "Acción (a_t) y Función de Recompensa (r_t)", font_size=20, bold=True, color=ACCENT2)
add_para(tf, "El Espacio de Acciones:", font_size=15, color=WHITE, bold=True)
add_para(tf, "Es el peso continuo asignado a cada activo. Se usa una capa Softmax para asegurar que w_i >= 0 y Suman 1.", font_size=13, color=GRAY)
add_para(tf, "", font_size=5)
add_para(tf, "Recompensa Media-Varianza (Igual que GAMS):", font_size=15, color=WHITE, bold=True)
add_para(tf, "r_t = w^T μ_mix - λ w^T Σ_mix w - Costos", font_size=15, color=ORANGE, font_name="Courier New", bold=True)
add_para(tf, "• Se utilizará λ = 0.10, correspondiente al 'Escenario 2' del modelo GAMS.", font_size=13, color=GRAY)
add_para(tf, "", font_size=5)
add_para(tf, "Dinámica (Capital y Costos):", font_size=15, color=WHITE, bold=True)
add_para(tf, "Costos asimétricos: 0.5% SPX, 1.0% CMC200.", font_size=13, color=GRAY)
add_para(tf, "x_{t+1} = x_t (1 + w_t^T R_t - Σ_i c_i |w_{i,t} - w_{i,t-1}|)", font_size=13, color=ACCENT, font_name="Courier New")

# ════════════════════════════════════════════════════════════════
# SLIDE 4 – BAJO EL CAPÓ: PPO
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.5, 11, 0.6, "3. Bajo el Capó de PPO: Proximal Policy Optimization", font_size=32, bold=True, color=ACCENT)

# Proceso On Policy
add_card(s, 0.8, 1.5, 11.7, 1.4)
tf = add_text(s, 1.1, 1.7, 11.1, 0.5, "Recolección (On-Policy)", font_size=18, bold=True, color=ORANGE)
add_para(tf, "El agente realiza Rollouts (ej. 2052 pasos, múltiplo exacto de 114 semanas) con la política congelada.", font_size=14, color=WHITE)
add_para(tf, "La red no entrega porcentajes directos, entrega una Campana de Gauss (μ, σ) de la cual se muestrean acciones.", font_size=14, color=GRAY)

add_card(s, 0.8, 3.2, 5.6, 3.5)
tf = add_text(s, 1.1, 3.5, 5, 0.5, "El Crítico y la Función Advantage", font_size=18, bold=True, color=ACCENT2)
add_para(tf, "El Crítico predice V(S_t) (cuánto ganaremos).", font_size=14, color=WHITE, bold=True)
add_para(tf, "Ventaja (A_t) = Retorno Real - Predicción V(S_t)", font_size=14, color=ACCENT, font_name="Courier New")
add_para(tf, "Si A_t > 0: La decisión fue inesperadamente buena.", font_size=14, color=GRAY)
add_para(tf, "Si A_t < 0: La decisión fue mala.", font_size=14, color=GRAY)

add_card(s, 6.7, 3.2, 5.8, 3.5)
tf = add_text(s, 7.0, 3.5, 5.2, 0.5, "Actualización: El Clipping y Optimizador", font_size=18, bold=True, color=PURPLE)
add_para(tf, "PPO usa ADAM en mini-batches.", font_size=14, color=WHITE, bold=True)
add_para(tf, "Clipping: PPO recorta las actualizaciones (max 20% cambio) para evitar que la red altere sus pesos bruscamente.", font_size=14, color=GRAY)
add_para(tf, "Garantía Constante: Entrenamiento altamente estable (no sufre caídas catastróficas del rendimiento).", font_size=14, color=LIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 – BAJO EL CAPÓ: SAC
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_text(s, 0.8, 0.5, 11, 0.6, "4. Bajo el Capó de SAC: Soft Actor-Critic", font_size=32, bold=True, color=ACCENT)

# Proceso Off Policy
add_card(s, 0.8, 1.5, 11.7, 1.4)
tf = add_text(s, 1.1, 1.7, 11.1, 0.5, "Replay Buffer y Ecosistema de 5 Redes (Off-Policy)", font_size=18, bold=True, color=ORANGE)
add_para(tf, "SAC guarda la experiencia pasada en un Replay Buffer y aprende en cada paso tomando muestras aleatorias (mini-batches).", font_size=14, color=WHITE)
add_para(tf, "Utiliza 1 Actor, 2 Críticos Gemelos y 2 Redes Target (rezagadas para dar estabilidad).", font_size=14, color=GRAY)

add_card(s, 0.8, 3.2, 5.6, 3.5)
tf = add_text(s, 1.1, 3.5, 5, 0.5, "Pesimismo y Críticos Q(s, a)", font_size=18, bold=True, color=ACCENT2)
add_para(tf, "El Truco del Pesimismo:", font_size=15, color=WHITE, bold=True)
add_para(tf, "El Crítico evalúa el valor Q(s,a) de un estado y acción.", font_size=14, color=GRAY)
add_para(tf, "SAC toma el valor mínimo entre el Crítico 1 y 2 para evitar el Sobre-Optimismo (Overestimation Bias) clásico en Q-learning continuo.", font_size=14, color=GRAY)

add_card(s, 6.7, 3.2, 5.8, 3.5)
tf = add_text(s, 7.0, 3.5, 5.2, 0.5, "Exploración: Máxima Entropía y Tau (τ)", font_size=18, bold=True, color=PURPLE)
add_para(tf, "Bono de Entropía (α x H):", font_size=15, color=WHITE, bold=True)
add_para(tf, "Recompensa a la red por mantener su campana de Gauss 'ancha' (explorar) vs casarse con una sola decisión tempranamente.", font_size=14, color=GRAY)
add_para(tf, "Suavizado Tau (τ = 0.005):", font_size=15, color=WHITE, bold=True)
add_para(tf, "Actualiza las redes Target solo en un 0.5% por paso, aportando alta estabilidad a la volatilidad del mercado.", font_size=14, color=GRAY)

# Guardar
out_path = os.path.join(os.path.dirname(__file__), "Presentacion_Profesor_Guia.pptx")
prs.save(out_path)
print(f"OK - Presentacion actualizada con el mega-resumen de la documentacion teórica.")
